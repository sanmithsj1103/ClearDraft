import logging
import json
import io
from fastapi import FastAPI, UploadFile, File, HTTPException, status
from fastapi.responses import StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, Field
from typing import Optional, List
from pptx import Presentation
from app.services.gemini import transcribe_audio, generate_polished_content, is_mock_active

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = FastAPI(
    title="ClearDraft API",
    description="Backend service for transcribing voice and generating structured formats using Gemini API",
    version="1.0"
)

# Enable CORS for frontend development server
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # For development, allow all origins. Can narrow to specific host in prod.
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# 25 MB file size limit
MAX_FILE_SIZE = 25 * 1024 * 1024

# Supported audio MIME types
SUPPORTED_MIME_TYPES = {
    "audio/mpeg", 
    "audio/mp3", 
    "audio/wav", 
    "audio/wave", 
    "audio/x-wav", 
    "audio/webm", 
    "audio/ogg",
    "audio/x-m4a",
    "audio/m4a",
    "audio/mp4"
}

class GenerateRequest(BaseModel):
    text: str = Field(..., min_length=1, description="The raw, unstructured thoughts or transcript")
    mode: str = Field(..., description="The target output format mode")
    tone: Optional[str] = Field(None, description="The tone setting (specifically for emails)")

# Slide formatting schemas for PowerPoint download
class SlideItem(BaseModel):
    title: str = Field(..., description="The title of the slide")
    type: str = Field(..., description="The type of the slide: 'title' or 'content'")
    subtitle: Optional[str] = Field(None, description="The subtitle (applies to title slide)")
    bullets: Optional[List[str]] = Field(None, description="List of bullet points (applies to content slides)")

class DownloadPPTRequest(BaseModel):
    slides: List[SlideItem] = Field(..., description="Array of slide details to compile")

@app.get("/health")
def health_check():
    """
    Checks backend health and API configuration status.
    """
    return {
        "status": "healthy",
        "mock_mode": is_mock_active(),
        "message": "ClearDraft backend is running smoothly!"
    }

@app.post("/api/transcribe")
async def transcribe(file: UploadFile = File(...)):
    """
    Receives an audio file, checks size constraints, and sends it to the Gemini API for transcription.
    """
    logger.info(f"Received transcription request for file: {file.filename}")
    
    # 1. Validate file size
    file_bytes = await file.read()
    file_size = len(file_bytes)
    
    if file_size > MAX_FILE_SIZE:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"File size exceeds maximum limit of 25MB (got {file_size / (1024 * 1024):.2f}MB)"
        )
        
    # 2. Validate MIME type
    mime_type = file.content_type
    
    if mime_type == "application/octet-stream" or not mime_type:
        filename_lower = file.filename.lower()
        if filename_lower.endswith(".mp3"):
            mime_type = "audio/mpeg"
        elif filename_lower.endswith(".wav"):
            mime_type = "audio/wav"
        elif filename_lower.endswith(".webm"):
            mime_type = "audio/webm"
        elif filename_lower.endswith(".m4a"):
            mime_type = "audio/x-m4a"
        elif filename_lower.endswith(".ogg"):
            mime_type = "audio/ogg"
            
    logger.info(f"Processing audio with content-type: {mime_type}")
    
    # 3. Transcribe audio using Gemini
    try:
        transcript = transcribe_audio(file_bytes, mime_type)
        return {
            "success": True,
            "filename": file.filename,
            "transcript": transcript
        }
    except Exception as e:
        logger.error(f"Failed to transcribe audio file {file.filename}: {str(e)}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Transcription failed: {str(e)}"
        )

@app.post("/api/generate")
def generate(payload: GenerateRequest):
    """
    Polishes raw transcripts/thoughts into the selected output mode using Gemini prompts.
    """
    logger.info(f"Received generation request for mode: {payload.mode}")
    
    try:
        polished = generate_polished_content(
            text=payload.text,
            mode=payload.mode.lower(),
            tone=payload.tone
        )
        
        # If mode is PPT, parse JSON and format a textual preview response
        if payload.mode.lower() == "ppt":
            try:
                slides_data = json.loads(polished)
                preview_parts = []
                for idx, slide in enumerate(slides_data, 1):
                    preview_parts.append(f"--- Slide {idx}: {slide.get('title', 'Untitled')} ---")
                    if slide.get('type') == 'title':
                        if 'subtitle' in slide and slide['subtitle']:
                            preview_parts.append(f"Subtitle: {slide['subtitle']}")
                    else:
                        bullets = slide.get('bullets', [])
                        for bullet in bullets:
                            preview_parts.append(f"• {bullet}")
                    preview_parts.append("")  # spacer line
                
                preview_text = "\n".join(preview_parts)
                return {
                    "success": True,
                    "mode": payload.mode,
                    "output": preview_text.strip(),
                    "slides": slides_data
                }
            except json.JSONDecodeError as jde:
                logger.error(f"Failed to parse Gemini output as JSON: {polished}. Error: {str(jde)}")
                raise HTTPException(
                    status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                    detail="Gemini API returned an invalid presentation outline structure. Please try again."
                )

        return {
            "success": True,
            "mode": payload.mode,
            "output": polished
        }
    except Exception as e:
        logger.error(f"Failed to generate polished output: {str(e)}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Content generation failed: {str(e)}"
        )

@app.post("/api/download-ppt")
def download_ppt(payload: DownloadPPTRequest):
    """
    Constructs a real PPTX presentation in-memory and returns the binary stream for download.
    """
    logger.info(f"Compiling PPTX presentation from {len(payload.slides)} slides data.")
    
    try:
        prs = Presentation()
        
        for slide_data in payload.slides:
            if slide_data.type == "title":
                # Title slide layout (0)
                slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(slide_layout)
                
                # Title
                if slide.shapes.title:
                    slide.shapes.title.text = slide_data.title
                # Subtitle (placeholder 1)
                if len(slide.placeholders) > 1 and slide_data.subtitle:
                    slide.placeholders[1].text = slide_data.subtitle
            else:
                # Content slide layout (1)
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                
                # Title
                if slide.shapes.title:
                    slide.shapes.title.text = slide_data.title
                # Bullets (placeholder 1 text frame)
                if len(slide.placeholders) > 1 and slide_data.bullets:
                    tf = slide.placeholders[1].text_frame
                    tf.text = ""
                    for idx, bullet in enumerate(slide_data.bullets):
                        if idx == 0:
                            tf.text = bullet
                        else:
                            p = tf.add_paragraph()
                            p.text = bullet
        
        # Save PPTX into in-memory bytes stream
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        
        return StreamingResponse(
            buffer,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": "attachment; filename=presentation.pptx"}
        )
    except Exception as e:
        logger.error(f"Failed to assemble PPTX file: {str(e)}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"PowerPoint generation failed: {str(e)}"
        )

if __name__ == "__main__":
    import uvicorn
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=True)
